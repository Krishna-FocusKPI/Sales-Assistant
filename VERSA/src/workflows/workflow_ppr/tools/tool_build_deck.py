import copy
import logging
import math
import os
import uuid

import streamlit as st
try:
    from langchain_core.tools import BaseTool
except ImportError:
    from langchain.tools import BaseTool

from src.common.workflow_context import get_workflow
from src.utils.versa_paths import get_versa_downloads_dir, resolve_workflow_path
from PIL import Image


def _get_deck_config(key: str, secrets_dot_path: str):
    """Get deck config from st.secrets (when in main thread) or os.environ (worker thread)."""
    try:
        obj = st.secrets
        for part in secrets_dot_path.split("."):
            obj = obj[part]
        if obj is not None:
            return str(obj).strip()
    except (AttributeError, KeyError, TypeError, Exception):
        pass
    env_map = {
        "template_path": "PPR_DECK_TEMPLATE_PATH",
        "item_per_listing_page": "PPR_DECK_ITEM_PER_LISTING_PAGE",
        "deck_saving_path": "PPR_DECK_SAVING_PATH",
        "page_intro": "PPR_DECK_PAGE_INTRO",
        "page_product_listing": "PPR_DECK_PAGE_PRODUCT_LISTING",
        "page_product_details": "PPR_DECK_PAGE_PRODUCT_DETAILS",
        "page_logo_sales_analysis": "PPR_DECK_PAGE_LOGO_SALES_ANALYSIS",
        "page_news_analysis": "PPR_DECK_PAGE_NEWS_ANALYSIS",
        "template_page_number": "PPR_DECK_TEMPLATE_PAGE_NUMBER",
    }
    return os.environ.get(env_map.get(key, f"PPR_DECK_{key.upper()}"), "").strip() or None
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_VERTICAL_ANCHOR as MSO_ANCHOR
from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT as PP_ALIGN
from pptx.util import Inches, Pt


class BuildDeck(BaseTool):
    name: str = "build_deck"
    description: str = (
        "Build the recurring revenue PowerPoint deck from workflow state. "
        "Requires analyze_logo_sales to have completed successfully first (non-empty analysis in memory). "
        "No parameters."
    )

    def _run(self) -> str:
        workflow = get_workflow()
        if not workflow:
            return "Error: workflow context not available."
        memory = workflow["workflow_memory"]
        logo_sales_analysis = getattr(memory, "logo_sales_analysis", None)
        if logo_sales_analysis is None or (
            hasattr(logo_sales_analysis, "empty") and logo_sales_analysis.empty
        ):
            return (
                "Cannot build the deck yet: logo sales analysis is missing or empty. "
                "Call analyze_logo_sales first (after distributor and logo are validated), "
                "then call build_deck when the user is ready for the presentation."
            )
        template_path = _get_deck_config("template_path", "ppr.deck.template_path")
        if not template_path:
            return "Error: deck template path not configured (set in secrets or PPR_DECK_TEMPLATE_PATH)."
        template_path = resolve_workflow_path(template_path)
        if not template_path or not os.path.isfile(template_path):
            return f"Error: deck template file not found at {template_path or 'resolved path'}."
        deck = Presentation(template_path)
        try:
            logging.info(f"* Building deck:")
            shopping_list = memory.shopping_list
            logo_name = memory.logo_name
            distributor_name = memory.distributor_name
            category = memory.category
            logo_sales_analysis_date = memory.logo_sales_analysis_date
            logo_sales_analysis = memory.logo_sales_analysis
            
            products_per_page = int(_get_deck_config("item_per_listing_page", "ppr.deck.item_per_listing_page") or "5")
            # Step 1: setup the intro page
            _page_intro(deck=deck)

            # Step 2: setup logo sales analysis page
            _page_logo_sales_analysis(
                deck=deck, 
                logo_name=logo_name, 
                distributor_name=distributor_name, 
                logo_sales_analysis_date=logo_sales_analysis_date, 
                logo_sales_analysis=logo_sales_analysis
            )
            
            # Step 2: setup the product listing page
            total_products = shopping_list.shape[0]
            total_pages = math.ceil(total_products / products_per_page)
            for page_number in range(total_pages):
                # get the start and end index for the page
                start_idx = page_number * products_per_page
                end_idx = min(start_idx + products_per_page, total_products)
                
                # get the products for the page
                page_products = shopping_list.iloc[start_idx:end_idx]
                
                # create the page
                _page_product_listing(
                    deck=deck, 
                    shopping_list=page_products,
                    logo_name=logo_name,
                    distributor_name=distributor_name,
                    category=category,
                    # yoy_analysis=yoy_analysis
                )
            
            # Step 3: setup the product details page
            for _, product in shopping_list.iterrows():
                _page_product_details(deck=deck, product=product)

            # Step 4: setup news analysis page
            if memory.need_news_analysis:
                for topic, news in memory.news_analysis.items():
                    _page_news_analysis(deck=deck, logo_name=logo_name, topic=topic, news=news)   
            
            _remove_template_slide(deck)

            filename = f"presentation_{uuid.uuid4()}.pptx"
            save_path = _get_deck_config("deck_saving_path", "downloads.deck_saving_path")
            if not save_path:
                return "Error: deck save path not configured (set in secrets or PPR_DECK_SAVING_PATH)."
            # Use project .versa/downloads when configured path doesn't exist (e.g. /downloads on Windows)
            save_path = save_path.replace("\\", "/").rstrip("/")
            if not os.path.isdir(save_path):
                fallback = get_versa_downloads_dir()
                save_path = str(fallback)
            os.makedirs(save_path, exist_ok=True)
            save_dir = os.path.join(save_path, filename)
            deck.save(save_dir)
            self._on_success(filename, save_dir)
            
        except Exception as e:
            logging.error(f"Error while building deck.", exc_info=True)
            return self._on_error()
        
    def _arun(self):
        raise NotImplementedError("This tool does not support async")
    
    def _on_success(self, filename: str, save_dir: str) -> str:
        workflow = get_workflow()
        if not workflow:
            return "Sucessfully built the deck."
        memory = workflow['workflow_memory']
        memory.deck_name = filename
        memory.deck_path = save_dir
        to_next = workflow['to_next_memory']
        to_next.reset()
        to_next.decision = "SUCCESS"
        to_next.source = self.name
        to_next.message = (
            "Successfully built the deck. Use **Download PowerPoint deck** under **Deck file** in the left sidebar "
            "to save the file to your computer."
        )
        logging.info(f"* Sucessfully built the deck at {save_dir}")
        return to_next.message

    def _on_error(self) -> str:
        workflow = get_workflow()
        if not workflow:
            return "An error occurred while building deck."
        to_next = workflow['to_next_memory']
        to_next.reset()
        to_next.decision = "ERROR"
        to_next.source = self.name
        to_next.message = f"An error occurred while building deck."
        return to_next.message


def _page_intro(deck):
    page_intro_idx = int(_get_deck_config("page_intro", "ppr.deck.page_intro") or "0")
    template_page = deck.slides[page_intro_idx]
    
    _add_new_slide(src_slide=template_page, deck=deck)


def _page_product_listing(deck, shopping_list, logo_name, distributor_name, category, yoy_analysis=None):
    def _set_title(shape):
        shape.text = shape.text.replace('[LOGO]', logo_name)
        
        # set font size and bold
        for paragraph in shape.text_frame.paragraphs:
            paragraph.alignment = PP_ALIGN.LEFT
            for run in paragraph.runs:
                run.font.size = Pt(20)
                run.font.bold = True
    
    def _set_subtitle(shape):
        # if yoy_analysis:
        #     shape.text = ''
        #     # p1
        #     paragraph = shape.text_frame.paragraphs[0]
        #     run = paragraph.add_run()
        #     run.text = f'Over the past year, {distributor_name.upper()} recurring revenue sales of {category.upper()} to {logo_name.upper()} have changed by {yoy_analysis}.'
        # else:
        #     shape.text = ''
        #     # p1
        #     paragraph = shape.text_frame.paragraphs[0]
        #     run = paragraph.add_run()
        #     run.text = f'Data not available for a full year-over-year analysis due to limited historical recurring sales data with {logo_name.upper()}.'
        shape.text = shape.text.replace('[LOGO]', logo_name)
        
        for paragraph in shape.text_frame.paragraphs:
            paragraph.alignment = PP_ALIGN.LEFT
            for run in paragraph.runs:
                run.font.size = Pt(12)
                run.font.bold = True
                run.font.italic = False

    def _set_table(shape):
        table = shape.table
        
        # expand the table since it only has one row
                    
        for _ in range(shopping_list.shape[0]-1):
            new_row = copy.deepcopy(table._tbl.tr_lst[-1])
            table._tbl.append(new_row)

        for i, row in enumerate(table.rows):
            # skip the header row
            if i == 0:
                continue
            
            cur_product = shopping_list.iloc[i - 1]

            # set item id
            row.cells[0].text_frame.text_anchor = MSO_ANCHOR.MIDDLE
            _run_1 = row.cells[0].text_frame.add_paragraph()
            _run_1.text = cur_product['ITEM_ID'].strip()
            _run_1.font.size = Pt(12)
            
            # set item name
            row.cells[1].text_frame.text_anchor = MSO_ANCHOR.MIDDLE
            _run_2 = row.cells[1].text_frame.add_paragraph()
            _run_2.text = cur_product['ITEM_NAME'].strip()
            _run_2.font.size = Pt(12)

            # set reason
            row.cells[2].text_frame.text_anchor = MSO_ANCHOR.MIDDLE
            reason_cell = row.cells[2]
            if reason_cell.text_frame.paragraphs:
                # Get the first paragraph in the cell
                paragraph = reason_cell.text_frame.paragraphs[0]
            else:
                # Add a new paragraph if none exists
                paragraph = reason_cell.text_frame.add_paragraph()

            run = paragraph.add_run()
            run.text = cur_product['REASON']
            run.font.size = Pt(12)
    
    page_product_listing_idx = int(_get_deck_config("page_product_listing", "ppr.deck.page_product_listing") or "1")
    template_page = deck.slides[page_product_listing_idx]
    page = _add_new_slide(template_page, deck)
    
    shape_dict = {shape.name: shape for shape in page.shapes}

    _set_title(shape_dict['title'])
    _set_subtitle(shape_dict['subtitle'])
    _set_table(shape_dict['table'])

def _page_product_details(deck, product):
    def _set_font_size(_cell, _title, _body, url=None):
        _cell.text = ''
        _title = str(_title)
        _body = str(_body) if _body else 'N/A'

        if _cell.text_frame.paragraphs:
            # Get the first paragraph in the cell
            _paragraph = _cell.text_frame.paragraphs[0]
        else:
            # Add a new paragraph if none exists
            _paragraph = _cell.text_frame.add_paragraph()

        _run = _paragraph.add_run()
        _run.text = f"{_title}\n"
        _run.font.size = Pt(10)
        _run.font.bold = True

        _run = _paragraph.add_run()
        _run.text = f"{_body}"
        _run.font.size = Pt(10)
        if url:
            link = _run.hyperlink
            link.address = url
                
    def _set_table_cell(cell):
        pricing = product['QUANTITY_BREAK_PRICING']
        
        if cell.text == '[ITEM NAME]':
            _set_font_size(cell, "ITEM NAME", product['ITEM_NAME'], product['URL'])
        elif cell.text == '[ITEM ID]':
            _set_font_size(cell, "ITEM ID", product['ITEM_ID'])
        elif cell.text == '[IS_PROUD_PATH]':
            _set_font_size(cell, "PROUD PATH", 'Yes' if product['IS_PROUD_PATH'] is True else 'No')
        elif cell.text == '[DESCRIPTION]':
            _set_font_size(cell, "DESCRIPTION", product['DESCRIPTION'])
        elif cell.text == '[HIGHLIGHTS]':
            _set_font_size(cell, "HIGHLIGHTS", product['HIGHLIGHTS'])
        elif cell.text == '[MOQ]':
            _set_font_size(cell, "MOQ", product['MOQ'])
        elif cell.text == '[BRAND]':
            _set_font_size(cell, "BRAND", product['BRAND'])
        elif f"{cell.text}".strip().replace('[', '').replace(']', '') in pricing:
            key = cell.text.strip().replace('[', '').replace(']', '')
                            
            cell.text = ''
            paragraph = cell.text_frame.paragraphs[0]
            paragraph.alignment = PP_ALIGN.CENTER

            run = paragraph.add_run()
            if "mc_us" in key or "mc_ca" in key:
                if "mc_us" in key:
                    run.text = str(pricing[key]) + ' us$' 
                if "mc_ca" in key:
                    run.text = str(pricing[key]) + ' cdn$'
            else:
                run.text = str(pricing[key])
            run.font.size = Pt(10)
                            
            if 'us$' in run.text or 'cdn$' in run.text:
                run.font.color.rgb = RGBColor(0, 0, 0)  # RGB values for black
    
    def _set_image(page, position_info):
        image_path = product['IMAGE_URL']
        
        if image_path:
            # Extract details from the placement dictionary
            max_img_width = position_info['max_img_width']
            max_img_height = position_info['max_img_height']
            left = position_info['left']
            top = position_info['top']

            # Load and scale the image
            img = Image.open(image_path)
            img_width, img_height = img.size
            scale_factor = min(max_img_width / img_width, max_img_height / img_height) * 0.8
            scaled_width = img_width * scale_factor
            scaled_height = img_height * scale_factor
            
            # Calculate the top-left position to center the image within the reserved space
            image_left = left + (max_img_width - scaled_width) / 2
            image_top = top + (max_img_height - scaled_height) / 2
            
            # Insert the image on the page at the calculated position
            page.shapes.add_picture(image_path, image_left, image_top, width=scaled_width, height=scaled_height)
    
    def _get_image_placement_details(shape):
        """
        This function calculates the placement details for an image based on the position of a table on the given page.

        Args:
            page: Slide object from python-pptx representing a single slide.

        Returns:
            A dictionary containing 'max_img_width', 'max_img_height', 'left', and 'top' for image placement.
        """
        if shape is None:
            raise ValueError("No table found on the page.")

        # Define the margins and available space for the image
        left_margin = Inches(0.5)
        left = left_margin
        right = shape.left - left_margin
        top = shape.top
        bottom = shape.top + shape.height

        # Calculate the maximum dimensions for the image
        max_img_width = right - left
        max_img_height = bottom - top
        
        return {
            'max_img_width': max_img_width,
            'max_img_height': max_img_height,
            'left': left,
            'top': top
        }
         
    page_product_detail_idx = int(_get_deck_config("page_product_details", "ppr.deck.page_product_details") or "2")
    slide = deck.slides[page_product_detail_idx]
    page = _add_new_slide(slide, deck)

    # find the table shape
    table = None
    for shape in page.shapes:
        if shape.name == 'table':
            table = shape.table
            position_info = _get_image_placement_details(shape)
            break
    
    if table is None:
        raise ValueError("Table shape not found in the template page")
    
    # set the values in the table
    for row in table.rows:
        for cell in row.cells:
            _set_table_cell(cell)

    # set the image
    _set_image(page, position_info)


def _add_new_slide(src_slide, deck):
    layout = src_slide.slide_layout
    dst_slide = deck.slides.add_slide(layout)

    # remove auto added
    for shape in dst_slide.shapes:
        sp = shape.element
        sp.getparent().remove(sp)

    # Copy content from the slide you want to duplicate to the new slide
    for shape in src_slide.shapes:
        el = shape.element
        new_el = copy.deepcopy(el)
        dst_slide.shapes._spTree.insert_element_before(new_el, 'p:extLst')

    return dst_slide


def _remove_template_slide(deck):
    def delete_slide(presentation, slide_index):
        slide_id = presentation.slides._sldIdLst[slide_index].id
        # Check if the slide ID is in the relationships
        if slide_id in presentation.part.rels:
            presentation.part.drop_rel(slide_id)
        presentation.slides._sldIdLst.remove(presentation.slides._sldIdLst[slide_index])

    _n = int(_get_deck_config("template_page_number", "ppr.deck.template_page_number") or "3")
    for _ in range(_n):
        if len(deck.slides) > 0:  # Check if there are slides to delete
            delete_slide(deck, 0)
        else:
            break  # Exit the loop if there are no more slides to delete


def _page_logo_sales_analysis(deck, logo_name, distributor_name, logo_sales_analysis_date, logo_sales_analysis):
    def _set_title(shape):
        shape.text = shape.text.replace('[logo]', logo_name.upper())


    def _set_subtitle(shape):
        _dist = (distributor_name or "").upper()
        _logo = (logo_name or "").upper()
        _date = logo_sales_analysis_date if logo_sales_analysis_date is not None else ""
        shape.text = shape.text.replace("[DISTRIBUTOR]", _dist).replace("[LOGO]", _logo).replace("[DATE]", _date)
    
    def _set_table(shape):
        table = shape.table
        
        for i, row in enumerate(logo_sales_analysis.itertuples(), start=1):  # Start at 1 to skip the title row
            # Set index in the first column
            table.cell(i, 0).text = str(row.Index)
            run = table.cell(i, 0).text_frame.paragraphs[0].runs[0]
            run.font.size = Pt(11)

            # Fill other columns
            for j, value in enumerate(row[1:], start=1):  # row[1:] skips the index of the DataFrame
                cell = table.cell(i, j)
                cell.text = str(value)
                run = cell.text_frame.paragraphs[0].runs[0]
                run.font.size = Pt(11)

        # bold row
        for cell in table.rows[3].cells:  # Row indices are zero-based
            for paragraph in cell.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.bold = True

        # bold col
        for row in table.rows:
            cell = row.cells[8]  # Column indices are zero-based
            for paragraph in cell.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.bold = True

        for row in table.rows:
            for cell in row.cells:
                # Set vertical alignment to middle
                cell.text_frame.text_anchor = MSO_ANCHOR.MIDDLE

                # Set horizontal alignment to center for each paragraph in the cell
                for paragraph in cell.text_frame.paragraphs:
                    paragraph.alignment = PP_ALIGN.CENTER
                            
    page_logo_sales_analysis_idx = int(_get_deck_config("page_logo_sales_analysis", "ppr.deck.page_logo_sales_analysis") or "3")
    template_page = deck.slides[page_logo_sales_analysis_idx]
    page = _add_new_slide(template_page, deck)
    
    shape_dict = {shape.name: shape for shape in page.shapes}

    _set_title(shape_dict['title'])
    _set_subtitle(shape_dict['subtitle'])

    if logo_sales_analysis is None:
        logging.warning("logo_sales_analysis is None; skipping logo sales table (run logo sales analysis first).")
        return
    if hasattr(logo_sales_analysis, "empty") and logo_sales_analysis.empty:
        logging.warning("logo_sales_analysis is empty; skipping logo sales table.")
        return

    _set_table(shape_dict['table'])


def _page_news_analysis(deck, logo_name, topic, news):
    def _set_title(shape):
        shape.text = shape.text.replace('[LOGO]', logo_name.upper())
        
    def _set_news(shape):
        shape.text_frame.clear()
        run = shape.text_frame.add_paragraph().add_run()
        run.text = topic
        run.font.bold = True
        run.font.size = Pt(12)
        
        for idx, bullet in enumerate(news):
            if idx > 5:
                break
            p = shape.text_frame.add_paragraph()
            p.level = 1

            body_run = p.add_run()
            body_run.text = f"{bullet.date} - {bullet.summary} "
            body_run.font.size = Pt(12)

            more_run = p.add_run()
            more_run.text = "More"
            more_run.font.size = Pt(12)
            hlink = more_run.hyperlink
            hlink.address = bullet.url
        
            
    page_news_analysis_idx = int(_get_deck_config("page_news_analysis", "ppr.deck.page_news_analysis") or "4")
    template_page = deck.slides[page_news_analysis_idx]
    page = _add_new_slide(template_page, deck)
    
    shape_dict = {shape.name: shape for shape in page.shapes}

    _set_title(shape_dict['title'])
    _set_news(shape_dict['news'])