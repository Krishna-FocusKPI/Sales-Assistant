import copy
import logging
import math
import os
import uuid

import streamlit as st
try:
    from langchain_core.tools import BaseTool
except ImportError:
    from langchain.tools import BaseTool
from PIL import Image

from src.common.workflow_context import get_workflow
from src.utils.deck_ttl_cleanup import cleanup_expired_generated_decks
from src.utils.versa_paths import get_versa_downloads_dir, resolve_workflow_path
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_VERTICAL_ANCHOR as MSO_ANCHOR
from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT as PP_ALIGN
from pptx.util import Inches, Pt


class BuildDeck(BaseTool):
    name: str = "build_deck"
    description: str = "Use this tool to build a recurring revenue deck, to use this tool you do not need to provide any parameter."
    
    def _run(self) -> str:
        try:
            raw_template = (st.secrets.ipr.deck.template_path or "").strip()
        except (AttributeError, KeyError, TypeError):
            raw_template = ""
        if not raw_template:
            return "Error: deck template path not configured (set ipr.deck.template_path in secrets)."
        template_path = resolve_workflow_path(raw_template)
        if not template_path or not os.path.isfile(template_path):
            return f"Error: deck template file not found. Resolved path: {template_path or 'none'} (expected under project .versa/workflows/, e.g. ipr_template.pptx)."
        deck = Presentation(template_path)        

        try:
            logging.info(f"* Building deck:")
            workflow = get_workflow()
            if not workflow:
                return "Workflow context not available."
            industry = workflow['workflow_memory'].industry
            shopping_list = workflow['workflow_memory'].shopping_list
            products_per_page = st.secrets.ipr.deck.item_per_listing_page

            # Step 1: setup the intro page
            _page_intro(deck=deck)

            # Step 2: setup the product listing page
            total_products = shopping_list.shape[0]
            total_pages = math.ceil(total_products / products_per_page)
            for page_number in range(total_pages):
                # get the start and end index for the page
                start_idx = page_number * products_per_page
                end_idx = min(start_idx + products_per_page, total_products)
                
                # get the products for the page
                page_products = shopping_list.iloc[start_idx:end_idx]
                
                # create the page
                _page_product_listing(deck=deck, shopping_list=page_products, industry=industry)
            
            # Step 3: setup the product details page
            for _, product in shopping_list.iterrows():
                _page_product_details(deck=deck, product=product)

            _remove_template_slide(deck)

            filename = f"presentation_{uuid.uuid4()}.pptx"
            try:
                save_path = (st.secrets.downloads.deck_saving_path or "").strip().replace("\\", "/").rstrip("/")
            except (AttributeError, KeyError, TypeError):
                save_path = ""
            if not save_path or not os.path.isdir(save_path):
                save_path = str(get_versa_downloads_dir())
            os.makedirs(save_path, exist_ok=True)
            save_dir = os.path.join(save_path, filename)
            deck.save(save_dir)
            self._on_success(filename, save_dir)
            
        except Exception as e:
            logging.error(f"Error while building deck.", exc_info=True)
            return self._on_error()
        
    def _arun(self):
        raise NotImplementedError("This tool does not support async")
    
    def _on_success(self, filename: str, save_dir: str) -> str:
        workflow = get_workflow()
        if not workflow:
            return "Workflow context not available."
        memory = workflow['workflow_memory']
        memory.deck_name = filename
        memory.deck_path = save_dir
        try:
            cleanup_expired_generated_decks(force=True)
        except Exception:
            logging.getLogger(__name__).debug("Deck TTL cleanup after build failed", exc_info=True)

        to_next = workflow['to_next_memory']
        to_next.reset()
        to_next.decision = "SUCCESS"
        to_next.source = self.name
        to_next.message = (
            "Successfully built the deck. Use **Download PowerPoint deck** under **Deck file** in the left sidebar "
            "to save the file to your computer."
        )

        logging.info(f"* Sucessfully built the deck at {save_dir}")
        return to_next.message
    
    def _on_error(self) -> str:
        workflow = get_workflow()
        if not workflow:
            return "Workflow context not available."
        to_next = workflow['to_next_memory']
        to_next.reset()
        to_next.decision = "ERROR"
        to_next.source = self.name
        to_next.message = f"An error occurred while building deck."
        return to_next.message


def _page_intro(deck):
    # get the intro page
    page_intro_idx = st.secrets.ipr.deck.page_intro
    template_page = deck.slides[page_intro_idx]
    
    _add_new_slide(src_slide=template_page, deck=deck)


def _page_product_listing(deck, shopping_list, industry):
    def _set_title(shape, industry):
        shape.text = shape.text.replace('[INDUSTRY]', industry)
        
        # set font size and bold
        for paragraph in shape.text_frame.paragraphs:
            paragraph.alignment = PP_ALIGN.LEFT
            for run in paragraph.runs:
                run.font.size = Pt(20)
                run.font.bold = True
    
    def _set_subtitle(shape, industry):
        shape.text = shape.text.replace('[INDUSTRY]', industry)
                
    def _set_table(shape, products):
        table = shape.table
        
        # expand the table since it only has one row
                    
        for _ in range(products.shape[0]-1):
            new_row = copy.deepcopy(table._tbl.tr_lst[-1])
            table._tbl.append(new_row)

        for i, row in enumerate(table.rows):
            # skip the header row
            if i == 0:
                continue
            
            cur_product = products.iloc[i - 1]

            # set item id
            row.cells[0].text_frame.text_anchor = MSO_ANCHOR.MIDDLE
            _run_1 = row.cells[0].text_frame.add_paragraph()
            _run_1.text = cur_product['ITEM_ID'].strip()
            _run_1.font.size = Pt(12)
            
            # set item name
            row.cells[1].text_frame.text_anchor = MSO_ANCHOR.MIDDLE
            _run_2 = row.cells[1].text_frame.add_paragraph()
            _run_2.text = cur_product['ITEM_NAME'].strip()
            _run_2.font.size = Pt(12)

            # set reason
            row.cells[2].text_frame.text_anchor = MSO_ANCHOR.MIDDLE
            reason_cell = row.cells[2]
            if reason_cell.text_frame.paragraphs:
                # Get the first paragraph in the cell
                paragraph = reason_cell.text_frame.paragraphs[0]
            else:
                # Add a new paragraph if none exists
                paragraph = reason_cell.text_frame.add_paragraph()

            run = paragraph.add_run()
            run.text = cur_product['REASON']
            run.font.size = Pt(12)
    
    page_product_listing_idx = st.secrets.ipr.deck.page_product_listing
    template_page = deck.slides[page_product_listing_idx]
    page = _add_new_slide(template_page, deck)
    
    shape_dict = {shape.name: shape for shape in page.shapes}

    _set_title(shape_dict['title'], industry)
    _set_subtitle(shape_dict['subtitle'], industry)
    _set_table(shape_dict['table'], products=shopping_list)

def _page_product_details(deck, product):
    def _set_font_size(_cell, _title, _body, url=None):
        _cell.text = ''
        _title = str(_title)
        _body = str(_body) if _body else 'N/A'

        if _cell.text_frame.paragraphs:
            # Get the first paragraph in the cell
            _paragraph = _cell.text_frame.paragraphs[0]
        else:
            # Add a new paragraph if none exists
            _paragraph = _cell.text_frame.add_paragraph()

        _run = _paragraph.add_run()
        _run.text = f"{_title}\n"
        _run.font.size = Pt(10)
        _run.font.bold = True

        _run = _paragraph.add_run()
        _run.text = f"{_body}"
        _run.font.size = Pt(10)
        if url:
            link = _run.hyperlink
            link.address = url
                
    def _set_table_cell(cell):
        pricing = product['QUANTITY_BREAK_PRICING']
        
        if cell.text == '[ITEM NAME]':
            _set_font_size(cell, "ITEM NAME", product['ITEM_NAME'], product['URL'])
        elif cell.text == '[ITEM ID]':
            _set_font_size(cell, "ITEM ID", product['ITEM_ID'])
        elif cell.text == '[IS_PROUD_PATH]':
            _set_font_size(cell, "PROUD PATH", 'Yes' if product['IS_PROUD_PATH'] is True else 'No')
        elif cell.text == '[DESCRIPTION]':
            _set_font_size(cell, "DESCRIPTION", product['DESCRIPTION'])
        elif cell.text == '[HIGHLIGHTS]':
            _set_font_size(cell, "HIGHLIGHTS", product['HIGHLIGHTS'])
        elif cell.text == '[MOQ]':
            _set_font_size(cell, "MOQ", product['MOQ'])
        elif cell.text == '[BRAND]':
            _set_font_size(cell, "BRAND", product['BRAND'])
        elif f"{cell.text}".strip().replace('[', '').replace(']', '') in pricing:
            key = cell.text.strip().replace('[', '').replace(']', '')
                            
            cell.text = ''
            paragraph = cell.text_frame.paragraphs[0]
            paragraph.alignment = PP_ALIGN.CENTER

            run = paragraph.add_run()
            if "mc_us" in key or "mc_ca" in key:
                if "mc_us" in key:
                    run.text = str(pricing[key]) + ' us$' 
                if "mc_ca" in key:
                    run.text = str(pricing[key]) + ' cdn$'
            else:
                run.text = str(pricing[key])
            run.font.size = Pt(10)
                            
            if 'us$' in run.text or 'cdn$' in run.text:
                run.font.color.rgb = RGBColor(0, 0, 0)  # RGB values for black
    
    def _set_image(page, position_info):
        image_path = product['IMAGE_URL']
        
        if image_path:
            # Extract details from the placement dictionary
            max_img_width = position_info['max_img_width']
            max_img_height = position_info['max_img_height']
            left = position_info['left']
            top = position_info['top']

            # Load and scale the image
            img = Image.open(image_path)
            img_width, img_height = img.size
            scale_factor = min(max_img_width / img_width, max_img_height / img_height) * 0.8
            scaled_width = img_width * scale_factor
            scaled_height = img_height * scale_factor
            
            # Calculate the top-left position to center the image within the reserved space
            image_left = left + (max_img_width - scaled_width) / 2
            image_top = top + (max_img_height - scaled_height) / 2
            
            # Insert the image on the page at the calculated position
            page.shapes.add_picture(image_path, image_left, image_top, width=scaled_width, height=scaled_height)

    def _get_image_placement_details(shape):
        """
        This function calculates the placement details for an image based on the position of a table on the given page.

        Args:
            page: Slide object from python-pptx representing a single slide.

        Returns:
            A dictionary containing 'max_img_width', 'max_img_height', 'left', and 'top' for image placement.
        """
        if shape is None:
            raise ValueError("No table found on the page.")

        # Define the margins and available space for the image
        left_margin = Inches(0.5)
        left = left_margin
        right = shape.left - left_margin
        top = shape.top
        bottom = shape.top + shape.height

        # Calculate the maximum dimensions for the image
        max_img_width = right - left
        max_img_height = bottom - top

        return {
            'max_img_width': max_img_width,
            'max_img_height': max_img_height,
            'left': left,
            'top': top
        }
          
    page_product_detail_idx = st.secrets.ipr.deck.page_product_details
    slide = deck.slides[page_product_detail_idx]
    page = _add_new_slide(slide, deck)

    # find the table shape
    table = None
    for shape in page.shapes:
        if shape.name == 'table':
            table = shape.table
            position_info = _get_image_placement_details(shape)
            break
    
    if table is None:
        raise ValueError("Table shape not found in the template page")
    
    # set the values in the table
    for row in table.rows:
        for cell in row.cells:
            _set_table_cell(cell)

    # set the image
    _set_image(page, position_info)


def _add_new_slide(src_slide, deck):
    layout = src_slide.slide_layout
    dst_slide = deck.slides.add_slide(layout)

    # remove auto added
    for shape in dst_slide.shapes:
        sp = shape.element
        sp.getparent().remove(sp)

    # Copy content from the slide you want to duplicate to the new slide
    for shape in src_slide.shapes:
        el = shape.element
        new_el = copy.deepcopy(el)
        dst_slide.shapes._spTree.insert_element_before(new_el, 'p:extLst')

    return dst_slide


def _remove_template_slide(deck):
    def delete_slide(presentation, slide_index):
        slide_id = presentation.slides._sldIdLst[slide_index].id
        # Check if the slide ID is in the relationships
        if slide_id in presentation.part.rels:
            presentation.part.drop_rel(slide_id)
        presentation.slides._sldIdLst.remove(presentation.slides._sldIdLst[slide_index])

    for _ in range(st.secrets.ipr.deck.template_page_number):
        if len(deck.slides) > 0:  # Check if there are slides to delete
            delete_slide(deck, 0)
        else:
            break  # Exit the loop if there are no more slides to delete
